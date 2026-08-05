import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_deck(out_filename):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(15, 23, 42)
    SAFFRON = RGBColor(234, 88, 12)
    INDIGO = RGBColor(79, 70, 229)
    TEAL = RGBColor(13, 148, 136)
    BG_LIGHT = RGBColor(248, 250, 252)
    CARD_WHITE = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(15, 23, 42)
    TEXT_MUTED = RGBColor(100, 116, 139)
    BORDER_LIGHT = RGBColor(203, 213, 225)
    EMERALD = RGBColor(16, 185, 129)
    ACCENT_BG = RGBColor(238, 242, 255)

    blank_layout = prs.slide_layouts[6]

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_hdr(slide, title_text, category_text="HIDDENYATRA • MAJOR PROJECT VIVA"):
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
        hdr.fill.solid(); hdr.fill.fore_color.rgb = CARD_WHITE; hdr.line.color.rgb = BORDER_LIGHT

        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.09))
        line.fill.solid(); line.fill.fore_color.rgb = SAFFRON; line.line.fill.background()

        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.18), Inches(3.4), Inches(0.32))
        pill.fill.solid(); pill.fill.fore_color.rgb = ACCENT_BG; pill.line.color.rgb = INDIGO
        tf_p = pill.text_frame; tf_p.word_wrap = True
        p_p = tf_p.paragraphs[0]; p_p.text = category_text.upper()
        p_p.font.size = Pt(10); p_p.font.bold = True; p_p.font.color.rgb = INDIGO; p_p.font.name = "Segoe UI"
        p_p.alignment = PP_ALIGN.CENTER

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.52), Inches(11.5), Inches(0.55))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = title_text
        p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = NAVY; p.font.name = "Segoe UI"

    def add_notes(slide, notes_text):
        tf = slide.notes_slide.notes_text_frame
        tf.text = notes_text

    def create_base(title, category="SYSTEM DESIGNS & FEATURES"):
        s = prs.slides.add_slide(blank_layout)
        set_bg(s, BG_LIGHT)
        add_hdr(s, title, category)
        return s

    # SLIDE 1: Title
    s1 = prs.slides.add_slide(blank_layout); set_bg(s1, NAVY)
    l = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    l.fill.solid(); l.fill.fore_color.rgb = SAFFRON; l.line.fill.background()

    card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT

    tpill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.15), Inches(4.8), Inches(0.38))
    tpill.fill.solid(); tpill.fill.fore_color.rgb = ACCENT_BG; tpill.line.color.rgb = INDIGO
    p = tpill.text_frame.paragraphs[0]; p.text = "FINAL YEAR MAJOR PROJECT VIVA • B.TECH CSE"; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = INDIGO; p.font.name = "Segoe UI"; p.alignment = PP_ALIGN.CENTER

    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.65), Inches(10.9), Inches(1.3))
    tf = tbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "H I D D E N Y A T R A"; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = SAFFRON; p.font.name = "Segoe UI"
    p2 = tf.add_paragraph(); p2.text = "GIS-Enabled Modern Tourism Discovery & AI Trip Optimization Platform"; p2.font.size = Pt(20); p2.font.bold = True; p2.font.color.rgb = NAVY; p2.font.name = "Segoe UI"; p2.space_before = Pt(4)

    mcard = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.1), Inches(10.933), Inches(3.2))
    mcard.fill.solid(); mcard.fill.fore_color.rgb = BG_LIGHT; mcard.line.color.rgb = BORDER_LIGHT

    tb_m = s1.shapes.add_textbox(Inches(1.5), Inches(3.25), Inches(10.3), Inches(2.9))
    tf_m = tb_m.text_frame; tf_m.word_wrap = True
    meta = [
        ("Degree & Course:", "Bachelor of Technology (B.Tech) in Computer Science & Engineering"),
        ("Project Domain:", "Full-Stack Web Systems, Spatial GIS Mapping & Algorithmic Recommendation Engines"),
        ("Student Presenters:", "Final Year CSE Major Project Team (Roll No: CSE-2026)"),
        ("Project Supervisor:", "Guide / Department Professor Name"),
        ("Institution Name:", "Department of Computer Science & Engineering, University Name"),
        ("Tech Stack:", "Python 3.13 (Flask), MySQL 8.0 (PyMySQL PooledDB), Leaflet.js v1.9.4, ES6 Vanilla JS")
    ]
    for k, v in meta:
        p = tf_m.add_paragraph(); p.text = f"{k}  {v}"; p.font.size = Pt(13); p.font.name = "Segoe UI"; p.font.color.rgb = TEXT_DARK; p.space_before = Pt(4)
    add_notes(s1, "Good morning respected faculty. Today I present HiddenYatra, a GIS-enabled tourism discovery and AI itinerary planning platform.")

    # SLIDE 2: Problem Statement
    s2 = create_base("Problem Statement & Domain Need", "PROJECT NEED & MOTIVATION")
    c1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.6))
    c1.fill.solid(); c1.fill.fore_color.rgb = CARD_WHITE; c1.line.color.rgb = BORDER_LIGHT
    tb = s2.shapes.add_textbox(Inches(1.1), Inches(1.65), Inches(5.0), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = "🔴 Existing Industry Deficits"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = NAVY
    b1 = [
        "Commercial Over-Centralization: Portals focus heavily on top 2-3 cities, ignoring hundreds of regional heritage sites.",
        "Fragmented Information: Travelers struggle to find exact GPS coordinates, road access, and verified local stays.",
        "Lack of Spatial Tools: Static travel blogs lack interactive marker clustering, distance calculators, or day routing.",
        "Unmoderated Stale Entries: Unofficial sites suffer from broken image links, stale entries, and zero governance."
    ]
    for b in b1:
        p = tf.add_paragraph(); p.text = "• " + b; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(10)

    c2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.7), Inches(5.6))
    c2.fill.solid(); c2.fill.fore_color.rgb = CARD_WHITE; c2.line.color.rgb = BORDER_LIGHT
    tb2 = s2.shapes.add_textbox(Inches(7.1), Inches(1.65), Inches(5.1), Inches(5.1))
    tf2 = tb2.text_frame; tf2.word_wrap = True; p = tf2.paragraphs[0]; p.text = "🎯 Real-World Impact & Engineering Need"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = SAFFRON
    b2 = [
        "Economic Loss for Local Artisans: Homestays and local guides lose potential tourism revenue.",
        "Poor Tourist Experience: Travelers manage fragmented directions and impractical commute paths.",
        "Lack of Heritage Archives: Historic monuments fade away without structured digital archives.",
        "Engineering Mandate: Need for a unified, secure GIS platform with automated trip planning."
    ]
    for b in b2:
        p = tf2.add_paragraph(); p.text = "• " + b; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(10)
    add_notes(s2, "Traditional commercial portals neglect regional tourism. HiddenYatra bridges this gap using Leaflet GIS spatial mapping.")

    # SLIDE 3: Objectives
    s3 = create_base("Project Objectives & Scope", "PROJECT GOALS")
    objs = [
        ("1. Interactive GIS Engine", "Deploy Leaflet.js v1.9.4 with dynamic marker clustering, category color pins, live geolocation, and straight-line Haversine distance calculations.", INDIGO),
        ("2. Algorithmic Trip Generator", "Build a spatial clustering engine grouping nearby destinations into realistic 1 to 5 day itineraries tailored to budget tiers and traveler personas.", SAFFRON),
        ("3. Structured Cultural Archive", "Digitally catalog Bihar's 38 districts with curated historical narratives, local food specialties, transport links, and verified accommodations.", TEAL),
        ("4. Enterprise Admin Governance", "Engineer an 11-module Admin Control Center featuring PBKDF2-SHA256 security, CSRF protection, soft-delete recycle bins, and audit logging.", EMERALD)
    ]
    for idx, (title, desc, col) in enumerate(objs):
        x = Inches(0.8 + (idx % 2) * 5.9); y = Inches(1.5 + (idx // 2) * 2.8)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.5))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.18), Inches(2.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
        tb = s3.shapes.add_textbox(x + Inches(0.35), y + Inches(0.2), Inches(5.0), Inches(2.1))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(12); p2.font.color.rgb = TEXT_DARK; p2.space_before = Pt(6)
    add_notes(s3, "Our objectives focused on building a secure, high-performance web platform combining GIS mapping, trip generation, and governance.")

    # SLIDE 4: Proposed Solution Architecture
    s4 = create_base("Proposed Solution — HiddenYatra Platform", "SYSTEM SOLUTION")
    sol_boxes = [
        ("🗺️ GIS Mapping Engine", "• Leaflet.js v1.9.4\n• CARTO Vector Tiles\n• Marker Clustering\n• Live Geolocation", INDIGO),
        ("🤖 AI Trip Generator", "• Haversine Clustering\n• Persona Customizer\n• Budget Breakdown\n• One-Click PDF Export", SAFFRON),
        ("👥 Community Platform", "• User Suggestions\n• Visitor Photo Uploads\n• Ratings & Reviews\n• Wishlist Management", TEAL),
        ("🛡️ Governance & Security", "• 11 Admin Modules\n• PBKDF2 Password Hashing\n• Soft-Delete Recycle Bin\n• Immutable Audit Logs", EMERALD)
    ]
    for idx, (head, body, color) in enumerate(sol_boxes):
        x = Inches(0.8 + idx * 2.95); y = Inches(1.5)
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.75), Inches(5.4))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        hbar = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(0.85))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = color; hbar.line.fill.background()
        tf_h = hbar.text_frame; tf_h.word_wrap = True; p_h = tf_h.paragraphs[0]; p_h.text = head; p_h.font.size = Pt(13); p_h.font.bold = True; p_h.font.color.rgb = CARD_WHITE; p_h.alignment = PP_ALIGN.CENTER
        tb = s4.shapes.add_textbox(x + Inches(0.15), y + Inches(1.0), Inches(2.45), Inches(4.2))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = body; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(8)
    add_notes(s4, "HiddenYatra integrates spatial mapping, automated trip generation, user photo moderation, and full administrative governance.")

    # SLIDE 5: Project Overview & Stats
    s5 = create_base("Project Overview & Metrics", "EXECUTIVE SUMMARY")
    c1 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.6))
    c1.fill.solid(); c1.fill.fore_color.rgb = CARD_WHITE; c1.line.color.rgb = BORDER_LIGHT
    tb = s5.shapes.add_textbox(Inches(1.1), Inches(1.65), Inches(5.9), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = "📌 Core Project Attributes"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = NAVY
    pts = [
        ("Application Scope:", "Full-Stack GIS Tourism Portal for Bihar state."),
        ("Target Audience:", "Domestic/foreign tourists, backpackers, family planners."),
        ("Design System:", "Light/Dark mode theme engine, glassmorphism, grid layout."),
        ("Performance Standard:", "Sub-100ms API query response time via PooledDB reuse."),
        ("Quality Benchmark:", "87 unit & integration tests executed with 100% pass rate.")
    ]
    for k, v in pts:
        p = tf.add_paragraph(); p.text = f"• {k} {v}"; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(10)

    kpis = [("38", "Districts Cataloged", INDIGO), ("28+", "Deep-Dive Verified Spots", SAFFRON), ("11", "Admin Control Modules", TEAL), ("87", "Automated Tests Passed", EMERALD)]
    for idx, (num, lbl, col) in enumerate(kpis):
        x = Inches(7.6 + (idx % 2) * 2.5); y = Inches(1.6 + (idx // 2) * 2.7)
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(2.4))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        tb_k = s5.shapes.add_textbox(x, y + Inches(0.3), Inches(2.3), Inches(1.8))
        tf_k = tb_k.text_frame; tf_k.word_wrap = True; p_n = tf_k.paragraphs[0]; p_n.text = num; p_n.font.size = Pt(36); p_n.font.bold = True; p_n.font.color.rgb = col; p_n.alignment = PP_ALIGN.CENTER
        p_l = tf_k.add_paragraph(); p_l.text = lbl; p_l.font.size = Pt(11); p_l.font.bold = True; p_l.font.color.rgb = NAVY; p_l.alignment = PP_ALIGN.CENTER; p_l.space_before = Pt(4)
    add_notes(s5, "The platform features 38 districts, 28 detailed verified spots, 11 admin modules, and 87 automated unit tests.")

    # SLIDE 6: System Architecture Diagram
    s6 = create_base("System Architecture & Layered Design", "SYSTEM ARCHITECTURE")
    layers = [
        ("CLIENT LAYER", "Web Browser • HTML5 / CSS3 / Vanilla ES6 JS • Leaflet.js GIS Engine", INDIGO),
        ("APPLICATION LAYER (FLASK)", "Flask Blueprints (main, auth, places, itinerary, admin, api) • CSRF & PBKDF2 Engine", SAFFRON),
        ("DATABASE POOLING LAYER", "PyMySQL Connection Pool (dbutils.pooled_db.PooledDB) • Thread-Safe Pool (size=5, max=20)", TEAL),
        ("PERSISTENCE LAYER (MYSQL 8.0)", "Relational Database Engine • Foreign Key Cascades • Immutable Audit Logs • Soft Delete Bin", EMERALD)
    ]
    for idx, (title, desc, col) in enumerate(layers):
        y = Inches(1.5 + idx * 1.35)
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), y, Inches(11.333), Inches(1.15))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        badge = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y, Inches(2.4), Inches(1.15))
        badge.fill.solid(); badge.fill.fore_color.rgb = col; badge.line.fill.background()
        tf_b = badge.text_frame; tf_b.word_wrap = True; p_b = tf_b.paragraphs[0]; p_b.text = title; p_b.font.size = Pt(11); p_b.font.bold = True; p_b.font.color.rgb = CARD_WHITE; p_b.alignment = PP_ALIGN.CENTER
        tb = s6.shapes.add_textbox(Inches(3.6), y + Inches(0.2), Inches(8.5), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK
    add_notes(s6, "Our system architecture segregates responsibilities cleanly: Client layer, Flask application layer, database pooling layer, and MySQL persistence layer.")

    # SLIDE 7: Tech Stack
    s7 = create_base("Technology Stack & Justifications", "TECH STACK")
    table_shape = s7.shapes.add_table(7, 4, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
    table = table_shape.table
    table.columns[0].width = Inches(2.2); table.columns[1].width = Inches(2.2); table.columns[2].width = Inches(2.2); table.columns[3].width = Inches(5.133)
    headers = ["Layer", "Technology", "Version / Tool", "Justification & Purpose"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.text = h; p.font.bold = True; p.font.color.rgb = CARD_WHITE; p.font.size = Pt(12)
    data = [
        ("Frontend Core", "HTML5 / CSS3 / ES6 JS", "Native Standards", "Sub-millisecond DOM rendering, zero dependency bloat."),
        ("GIS & Maps", "Leaflet.js & CARTO", "v1.9.4 / Light Tiles", "High-contrast vector cartography with marker clustering."),
        ("Backend Framework", "Python & Flask", "v3.13 / v3.0+", "Lightweight WSGI micro-framework with Blueprint routing."),
        ("Database Engine", "MySQL", "v8.0+", "ACID compliance, relational schema, foreign key cascades."),
        ("Connection Pool", "PyMySQL + PooledDB", "DBUtils v3.0+", "Thread-safe pool (size=5, max=20) eliminating DB latency."),
        ("Security Hasher", "hashlib PBKDF2-SHA256", "Native Python", "Cryptographically secure password hashing with unique salts.")
    ]
    for r_idx, row in enumerate(data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx); cell.fill.solid(); cell.fill.fore_color.rgb = CARD_WHITE if r_idx % 2 == 0 else BG_LIGHT
            p = cell.text_frame.paragraphs[0]; p.text = val; p.font.size = Pt(11); p.font.color.rgb = TEXT_DARK
    add_notes(s7, "We selected Python and Flask for agility, MySQL for reliable relational queries, and Leaflet.js for responsive map interactions.")

    # SLIDE 8: Folder Structure
    s8 = create_base("Project Folder & Codebase Structure", "CODEBASE ARCHITECTURE")
    c1 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.6))
    c1.fill.solid(); c1.fill.fore_color.rgb = CARD_WHITE; c1.line.color.rgb = BORDER_LIGHT
    tb1 = s8.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(5.2), Inches(5.1))
    tf1 = tb1.text_frame; tf1.word_wrap = True; p = tf1.paragraphs[0]; p.text = "📂 Key Directory Modules"; p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = NAVY
    dirs = [
        ("app.py", "Flask application factory & server startup entrypoint."),
        ("models/database.py", "Unified database layer with PooledDB connection pool."),
        ("routes/", "Flask Blueprints (main, auth, places, itinerary, admin, api)."),
        ("static/", "CSS design system, map.js, app.js, dynamic file uploads."),
        ("templates/", "Jinja2 HTML5 templates, admin layout, base.html."),
        ("tests/", "Pytest/Unittest suite containing 87 test cases.")
    ]
    for d, desc in dirs:
        p = tf1.add_paragraph(); p.text = f"• {d}: {desc}"; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(8)

    c2 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.7), Inches(5.6))
    c2.fill.solid(); c2.fill.fore_color.rgb = CARD_WHITE; c2.line.color.rgb = BORDER_LIGHT
    tb2 = s8.shapes.add_textbox(Inches(7.0), Inches(1.65), Inches(5.3), Inches(5.1))
    tf2 = tb2.text_frame; tf2.word_wrap = True; p = tf2.paragraphs[0]; p.text = "⚙️ Blueprint Division"; p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = SAFFRON
    blueprints = [
        ("routes/main.py", "Homepage, State/District detail, Explore map routes."),
        ("routes/places.py", "Place details, Wishlist, Visited tracks, Reviews."),
        ("routes/itinerary.py", "AI Trip Planner & sequential route engine."),
        ("routes/auth.py", "User registration, login, PBKDF2 hashing, session."),
        ("routes/admin.py", "11 Admin control modules, approvals, audit logs."),
        ("routes/api.py", "REST endpoints (/api/autocomplete, /api/smart-search).")
    ]
    for bp, desc in blueprints:
        p = tf2.add_paragraph(); p.text = f"• {bp}: {desc}"; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(8)
    add_notes(s8, "Our codebase segregates responsibilities across modular Blueprints, template views, static assets, and automated unit tests.")

    # SLIDE 9: Functional Modules
    s9 = create_base("Functional Modules Breakdown", "SYSTEM MODULES")
    mods = [
        ("1. Public Discovery Portal", "Homepage hero media, 38 district catalogues, state overview pages, food & culture guides, and place search filters.", INDIGO),
        ("2. Interactive GIS Mapping", "Fullscreen GIS map, CARTO light tiles, dynamic marker clustering, category color pins, and live distance calculations.", SAFFRON),
        ("3. AI Trip Planner Engine", "Haversine coordinate clustering, 1-5 day itineraries, budget breakdowns (Budget/Standard/Luxury), and PDF export.", TEAL),
        ("4. Admin Control Center", "11 governance modules, place CRUD, user suggestion queue, visitor photo approvals, recycle bin, and audit logs.", EMERALD)
    ]
    for idx, (title, desc, col) in enumerate(mods):
        x = Inches(0.8 + (idx % 2) * 5.9); y = Inches(1.5 + (idx // 2) * 2.8)
        card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.5))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        hbar = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(0.55))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = col; hbar.line.fill.background()
        tf_h = hbar.text_frame; tf_h.word_wrap = True; p_h = tf_h.paragraphs[0]; p_h.text = title; p_h.font.size = Pt(14); p_h.font.bold = True; p_h.font.color.rgb = CARD_WHITE
        tb = s9.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), Inches(5.2), Inches(1.75))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK
    add_notes(s9, "HiddenYatra is divided into four major modules: Public Portal, GIS Mapping, AI Trip Planner, and Admin Control Center.")

    # SLIDE 10: DB Schema & ER
    s10 = create_base("Database Schema & ER Diagram", "DATABASE ARCHITECTURE")
    c1 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(4.5), Inches(5.6))
    c1.fill.solid(); c1.fill.fore_color.rgb = CARD_WHITE; c1.line.color.rgb = BORDER_LIGHT
    tb1 = s10.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(4.1), Inches(5.1))
    tf1 = tb1.text_frame; tf1.word_wrap = True; p = tf1.paragraphs[0]; p.text = "🗄️ Relational Schema Highlights"; p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = NAVY
    db_items = [
        "14 Relational Tables: users, states, districts, blocks, places, place_photos, specialties, accommodations, reviews, wishlist, user_visited, user_photos, itineraries, admin_logs.",
        "Foreign Keys & Cascades: Maintains integrity on deletion (e.g. district -> places -> photos).",
        "Soft Delete Support: Places table contains is_deleted flag for recycle bin recovery.",
        "Immutable Audit Logs: admin_logs records every admin action with IP & timestamps."
    ]
    for b in db_items:
        p = tf1.add_paragraph(); p.text = "• " + b; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(8)

    entities = [
        ("USERS", "id (PK)\nemail (UK)\npassword_hash\nfull_name\ncreated_at", INDIGO),
        ("PLACES", "id (PK)\ndistrict_id (FK)\nname, slug (UK)\ncategory, lat, lng\nis_deleted", SAFFRON),
        ("DISTRICTS", "id (PK)\nstate_id (FK)\nname, slug (UK)\nfamous_for\ncover_image", TEAL),
        ("ITINERARIES", "id (PK)\nuser_id (FK)\nname, days_count\nbudget_tier\ncreated_at", EMERALD)
    ]
    for idx, (ename, fields, col) in enumerate(entities):
        x = Inches(5.6 + (idx % 2) * 3.6); y = Inches(1.4 + (idx // 2) * 2.8)
        card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.3), Inches(2.5))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        hbar = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.3), Inches(0.5))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = col; hbar.line.fill.background()
        tf_h = hbar.text_frame; p_h = tf_h.paragraphs[0]; p_h.text = "TABLE: " + ename; p_h.font.size = Pt(12); p_h.font.bold = True; p_h.font.color.rgb = CARD_WHITE; p_h.alignment = PP_ALIGN.CENTER
        tb = s10.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55), Inches(3.0), Inches(1.85))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = fields; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK
    add_notes(s10, "Our database schema consists of 14 relational tables enforcing foreign keys and soft-delete recovery to prevent accidental data loss.")

    # SLIDE 11: Workflow
    s11 = create_base("End-to-End Workflow & User Journey", "WORKFLOW")
    steps = [
        ("1. Discovery & Search", "User searches destination or filters by district/category via /api/smart-search.", INDIGO),
        ("2. Interactive GIS View", "Leaflet GIS engine displays category pins, marker clusters, and straight-line distance.", SAFFRON),
        ("3. AI Trip Generation", "User specifies duration, persona, and budget tier. Haversine clusters create day-wise plans.", TEAL),
        ("4. Community Platform", "Registered users add reviews, mark visited spots, and submit new photos or place proposals.", EMERALD),
        ("5. PDF Export & Save", "User exports trip plan as clean print-optimized PDF or saves to personal profile itinerary.", NAVY)
    ]
    for idx, (title, desc, col) in enumerate(steps):
        x = Inches(0.8 + idx * 2.35); y = Inches(2.0)
        card = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(4.5))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        hbar = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.2), Inches(0.75))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = col; hbar.line.fill.background()
        tf_h = hbar.text_frame; tf_h.word_wrap = True; p_h = tf_h.paragraphs[0]; p_h.text = title; p_h.font.size = Pt(12); p_h.font.bold = True; p_h.font.color.rgb = CARD_WHITE; p_h.alignment = PP_ALIGN.CENTER
        tb = s11.shapes.add_textbox(x + Inches(0.1), y + Inches(0.9), Inches(2.0), Inches(3.4))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK
    add_notes(s11, "The user journey transitions smoothly from search and map exploration to automated itinerary generation and PDF export.")

    # SLIDE 12: GIS Exploration
    s12 = create_base("GIS Exploration & Leaflet Map Integration", "GIS ENGINE")
    c1 = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.6))
    c1.fill.solid(); c1.fill.fore_color.rgb = CARD_WHITE; c1.line.color.rgb = BORDER_LIGHT
    tb = s12.shapes.add_textbox(Inches(1.1), Inches(1.65), Inches(5.9), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = "🗺️ Leaflet GIS Features"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = NAVY
    gis_pts = [
        ("Leaflet.js v1.9.4 Engine:", "Lightweight open-source mapping engine with sub-10ms marker render time."),
        ("CARTO Light Tiles:", "High-contrast vector cartography tile layers configured in light mode."),
        ("Marker Clustering:", "Leaflet.markercluster automatically clusters dense destination markers into interactive count bubbles."),
        ("Category Pin Coding:", "Distinct color-coded pins (Temples=Yellow, Waterfalls=Blue, Forts=Purple, Nature=Green)."),
        ("Live Distance Calculation:", "Browser Geolocation API + Haversine formula estimates exact km distance from traveler.")
    ]
    for k, v in gis_pts:
        p = tf.add_paragraph(); p.text = f"• {k} {v}"; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(8)

    c2 = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(1.4), Inches(4.9), Inches(5.6))
    c2.fill.solid(); c2.fill.fore_color.rgb = NAVY; c2.line.color.rgb = INDIGO
    tb_c = s12.shapes.add_textbox(Inches(7.85), Inches(1.65), Inches(4.4), Inches(5.1))
    tf_c = tb_c.text_frame; tf_c.word_wrap = True; p = tf_c.paragraphs[0]; p.text = "💻 Haversine Formula (JS)"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = SAFFRON
    code = """function getDistanceKm(lat1, lon1, lat2, lon2) {
  const R = 6371; // Earth radius (km)
  const dLat = (lat2 - lat1) * Math.PI / 180;
  const dLon = (lon2 - lon1) * Math.PI / 180;
  const a = 
    Math.sin(dLat/2) * Math.sin(dLat/2) +
    Math.cos(lat1 * Math.PI / 180) * 
    Math.cos(lat2 * Math.PI / 180) * 
    Math.sin(dLon/2) * Math.sin(dLon/2);
  const c = 2 * Math.atan2(Math.sqrt(a), Math.sqrt(1-a));
  return R * c;
}"""
    p2 = tf_c.add_paragraph(); p2.text = code; p2.font.size = Pt(11); p2.font.name = "Consolas"; p2.font.color.rgb = CARD_WHITE; p2.space_before = Pt(10)
    add_notes(s12, "Leaflet.js provides high-performance spatial mapping, marker clustering, and instant distance calculations using the Haversine formula.")

    # SLIDE 13: AI Trip Planner
    s13 = create_base("AI Trip Planner & Itinerary Engine", "AI ROUTING PIPELINE")
    steps = [
        ("Step 1: User Preferences", "Inputs travel duration (1-5 days), persona (Solo/Family), and budget tier (Budget/Standard/Luxury).", INDIGO),
        ("Step 2: Spatial Selection", "Queries database for top destinations in target district matching preferred categories.", SAFFRON),
        ("Step 3: Haversine Clustering", "Groups spatial coordinates by geographic proximity to minimize daily commute times.", TEAL),
        ("Step 4: Route Polyline & PDF", "Draws Day 1 -> Day 2 connected polyline routes on map and generates exportable PDF.", EMERALD)
    ]
    for idx, (title, desc, col) in enumerate(steps):
        y = Inches(1.5 + idx * 1.35)
        card = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(1.15))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        badge = s13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(2.8), Inches(1.15))
        badge.fill.solid(); badge.fill.fore_color.rgb = col; badge.line.fill.background()
        tf_b = badge.text_frame; tf_b.word_wrap = True; p_b = tf_b.paragraphs[0]; p_b.text = title; p_b.font.size = Pt(12); p_b.font.bold = True; p_b.font.color.rgb = CARD_WHITE; p_b.alignment = PP_ALIGN.CENTER
        tb = s13.shapes.add_textbox(Inches(3.8), y + Inches(0.18), Inches(8.5), Inches(0.8))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK
    add_notes(s13, "Our trip engine algorithmically clusters destinations by spatial distance, generating optimized day-by-day itineraries and PDF exports.")

    # SLIDE 14: Security
    s14 = create_base("Authentication & Security Architecture", "SECURITY SHIELD")
    sec_features = [
        ("🔒 PBKDF2-SHA256 Hashing", "Passwords hashed with 100,000 iterations and unique 16-byte random salt per user.", INDIGO),
        ("🛡️ CSRF Double-Submit Tokens", "All POST forms require valid session-backed _csrf_token validation.", SAFFRON),
        ("⚡ Sliding Window Rate Limiter", "In-memory rate limiting blocks IP addresses exceeding 5 failed logins/min.", TEAL),
        ("🛡️ SQL Injection Prevention", "100% parameterized PyMySQL queries with _escape_like wildcard handling.", EMERALD),
        ("🍪 Session Hardening", "HTTPOnly=True, SameSite=Lax, and secure cookie properties enabled.", NAVY)
    ]
    for idx, (title, desc, col) in enumerate(sec_features):
        x = Inches(0.8 + (idx % 2) * 5.9) if idx < 4 else Inches(3.75)
        y = Inches(1.5 + (idx // 2) * 1.8)
        w = Inches(5.6) if idx < 4 else Inches(5.8)
        card = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(1.6))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        hbar = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.42))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = col; hbar.line.fill.background()
        tf_h = hbar.text_frame; p_h = tf_h.paragraphs[0]; p_h.text = title; p_h.font.size = Pt(12); p_h.font.bold = True; p_h.font.color.rgb = CARD_WHITE
        tb = s14.shapes.add_textbox(x + Inches(0.15), y + Inches(0.48), w - Inches(0.3), Inches(1.05))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK
    add_notes(s14, "Security is built into the core framework using PBKDF2-SHA256 password hashing, CSRF token validation, rate limiting, and parameterized SQL queries.")

    # SLIDE 15: Admin Center
    s15 = create_base("Enterprise Admin Control Center (11 Modules)", "ADMIN GOVERNANCE")
    admin_mods = [
        "1. Metrics Dashboard: Real-time count of total places, reviews, user photos, and wishlists.",
        "2. Places Management: Full CRUD interface to create, update, and soft-delete tourist spots.",
        "3. User Suggestions Queue: Approve or reject community place proposals.",
        "4. Visitor Photo Approvals: Review user photo uploads before public display.",
        "5. District & State Manager: Configure district descriptions, famous items, and custom order.",
        "6. Hero Media Manager: Customize homepage banner slides, title overlays, and active status.",
        "7. Appearance Customizer: Configure brand logo, theme colors, and auth banner styles.",
        "8. Recycle Bin (Soft Delete): Restore deleted records with 1-click or permanently delete.",
        "9. Immutable Audit Logs: Log every admin action with IP notes and timestamp history.",
        "10. Review Moderation: Delete inappropriate user reviews or spam entries.",
        "11. User Account Manager: Manage user registrations, roles, and access control."
    ]
    c = s15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.6))
    c.fill.solid(); c.fill.fore_color.rgb = CARD_WHITE; c.line.color.rgb = BORDER_LIGHT
    tb = s15.shapes.add_textbox(Inches(1.1), Inches(1.65), Inches(11.133), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = "🏛️ 11 Enterprise Governance Modules"; p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = NAVY
    for m in admin_mods:
        p = tf.add_paragraph(); p.text = "• " + m; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(4)
    add_notes(s15, "The Admin Control Center features 11 modules including place CRUD, user photo approval queues, soft-delete recycle bins, and immutable audit logs.")

    # SLIDE 16: REST APIs
    s16 = create_base("RESTful API Architecture & Endpoints", "API SPECIFICATIONS")
    table_shape = s16.shapes.add_table(6, 4, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3))
    table = table_shape.table
    table.columns[0].width = Inches(1.5); table.columns[1].width = Inches(3.2); table.columns[2].width = Inches(4.5); table.columns[3].width = Inches(2.533)
    headers = ["HTTP Method", "Endpoint Path", "Description", "Response Format"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]; p.text = h; p.font.bold = True; p.font.color.rgb = CARD_WHITE; p.font.size = Pt(13)
    apis = [
        ("GET", "/api/autocomplete?q=...", "Instant search dropdown suggestions.", "JSON: { suggestions: [...] }"),
        ("GET", "/api/smart-search?q=...", "Intent-based natural language search.", "JSON: { places: [...], count }"),
        ("GET", "/api/visited/<place_id>", "Check user visited status for place.", "JSON: { visited: bool }"),
        ("POST", "/api/visited/toggle", "Toggle place visited status.", "JSON: { status: 'visited' }"),
        ("GET", "/api/place/<id>/nearby", "Fetch hotels & food spots near place.", "JSON: { hotels: [], foods: [] }")
    ]
    for r_idx, row in enumerate(apis):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx); cell.fill.solid(); cell.fill.fore_color.rgb = CARD_WHITE if r_idx % 2 == 0 else BG_LIGHT
            p = cell.text_frame.paragraphs[0]; p.text = val; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK
    add_notes(s16, "Our REST APIs deliver fast JSON responses for autocomplete, smart search, visited tracking, and nearby accommodation fetching.")

    # SLIDE 17: Challenges
    s17 = create_base("Technical Challenges Faced & Solutions", "ENGINEERING CHALLENGES")
    challenges = [
        ("Map Marker Pin Overlap", "High density of spots in close proximity caused overlapping map pins.", "Integrated Leaflet.markercluster to dynamically group dense pins into interactive count bubbles.", INDIGO),
        ("Jinja2 Datetime Slicing", "PyMySQL datetime objects caused TypeError when subscripted in Jinja2 templates.", "Assigned explicit string variable {% set ts = log.created_at | string %} before slicing.", SAFFRON),
        ("Unpositioned Map Tile Bursting", "Leaflet container lost tile CSS when stylesheet loaded asynchronously.", "Included leaflet.css globally in base.html head and set position: relative; overflow: hidden;.", TEAL),
        ("SQL Connection Exhaustion", "Concurrent HTTP requests exhausted unclosed MySQL database connections.", "Implemented thread-safe PooledDB connection pool (size=5, max=20) with Python context managers.", EMERALD)
    ]
    for idx, (prob, cause, sol, col) in enumerate(challenges):
        y = Inches(1.5 + idx * 1.35)
        card = s17.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(1.15))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        badge = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(2.8), Inches(1.15))
        badge.fill.solid(); badge.fill.fore_color.rgb = col; badge.line.fill.background()
        tf_b = badge.text_frame; tf_b.word_wrap = True; p_b = tf_b.paragraphs[0]; p_b.text = prob; p_b.font.size = Pt(12); p_b.font.bold = True; p_b.font.color.rgb = CARD_WHITE; p_b.alignment = PP_ALIGN.CENTER
        tb = s17.shapes.add_textbox(Inches(3.8), y + Inches(0.15), Inches(8.5), Inches(0.85))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = f"Cause: {cause}"; p.font.size = Pt(11); p.font.color.rgb = TEXT_MUTED
        p2 = tf.add_paragraph(); p2.text = f"Solution: {sol}"; p2.font.size = Pt(12); p2.font.bold = True; p2.font.color.rgb = TEXT_DARK; p2.space_before = Pt(3)
    add_notes(s17, "During development, we systematically diagnosed and resolved critical challenges around spatial marker clustering, Jinja2 template formatting, map tile rendering, and connection pooling.")

    # SLIDE 18: Testing
    s18 = create_base("Testing, QA & Automated Validation", "TEST REPORT")
    kpis = [("87", "Automated Tests Executed", INDIGO), ("100%", "Pass Rate (0 Failures)", EMERALD), ("3.5s", "Test Suite Execution Time", SAFFRON)]
    for idx, (num, lbl, col) in enumerate(kpis):
        y = Inches(1.5 + idx * 1.8)
        card = s18.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(3.5), Inches(1.55))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        tb = card.text_frame; tb.word_wrap = True; p = tb.paragraphs[0]; p.text = num; p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = col; p.alignment = PP_ALIGN.CENTER
        p2 = tb.add_paragraph(); p2.text = lbl; p2.font.size = Pt(12); p2.font.bold = True; p2.font.color.rgb = NAVY; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(2)

    c2 = s18.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(1.5), Inches(7.8), Inches(5.3))
    c2.fill.solid(); c2.fill.fore_color.rgb = CARD_WHITE; c2.line.color.rgb = BORDER_LIGHT
    tb2 = s18.shapes.add_textbox(Inches(4.95), Inches(1.75), Inches(7.3), Inches(4.8))
    tf2 = tb2.text_frame; tf2.word_wrap = True; p = tf2.paragraphs[0]; p.text = "🧪 Key Test Modules Executed"; p.font.size = Pt(17); p.font.bold = True; p.font.color.rgb = NAVY
    t_items = [
        ("test_database.py:", "Validates PyMySQL connection pool creation, query execution, and soft-delete SQL logic."),
        ("test_security.py:", "Validates PBKDF2 hash generation, salt uniqueness, compare_digest verification, and CSRF token checks."),
        ("test_routes.py:", "Validates HTTP 200, 302, 404 response status codes across public, user, and admin Blueprint endpoints."),
        ("Manual Visual Inspection:", "Verified responsive layouts, dark/light mode toggle, and Leaflet map pins across browsers.")
    ]
    for title, desc in t_items:
        p = tf2.add_paragraph(); p.text = f"• {title} {desc}"; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(10)
    add_notes(s18, "We maintain an 87-case automated unit test suite covering database operations, security mechanisms, and route endpoints with 100% pass rate.")

    # SLIDE 19: Advantages
    s19 = create_base("Key System Advantages", "SYSTEM BENCHMARKS")
    advs = [
        ("⚡ High Performance & Speed", "Zero heavy JavaScript framework bundle; initial DOM render completes under 150ms.", INDIGO),
        ("🗺️ True Spatial Intelligence", "Provides GIS marker clustering, route polylines, and user geolocation distance in real time.", SAFFRON),
        ("🛡️ Data Recovery & Soft Deletes", "Recycle bin ensures deleted entries can be restored instantly without data loss.", TEAL),
        ("📱 Fully Responsive Design", "Flawless UI appearance across desktop, laptop, tablet, and smartphone browsers.", EMERALD),
        ("🚀 Production Ready", "Equipped with logging, security headers, rate limiting, and Docker container support.", NAVY)
    ]
    for idx, (title, desc, col) in enumerate(advs):
        y = Inches(1.5 + idx * 1.08)
        card = s19.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.733), Inches(0.95))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        hbar = s19.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(3.0), Inches(0.95))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = col; hbar.line.fill.background()
        tf_h = hbar.text_frame; tf_h.word_wrap = True; p_h = tf_h.paragraphs[0]; p_h.text = title; p_h.font.size = Pt(11); p_h.font.bold = True; p_h.font.color.rgb = CARD_WHITE; p_h.alignment = PP_ALIGN.CENTER
        tb = s19.shapes.add_textbox(Inches(4.0), y + Inches(0.1), Inches(8.3), Inches(0.75))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = desc; p.font.size = Pt(12); p.font.color.rgb = TEXT_DARK
    add_notes(s19, "HiddenYatra delivers high speed, robust security, and deep spatial functionality without relying on heavy third-party framework bundles.")

    # SLIDE 20: Future Scope
    s20 = create_base("Future Scope & Enhancement Roadmap", "ROADMAP")
    phases = [
        ("Phase 1 (Q3 2026)", "Native Mobile App", "Develop cross-platform iOS & Android mobile application using Flutter.", INDIGO),
        ("Phase 2 (Q4 2026)", "Offline Map Caching", "Implement Service Worker PWA offline tile caching for remote areas.", SAFFRON),
        ("Phase 3 (Q1 2027)", "Multi-Language (i18n)", "Add Hindi, Maithili, and Bhojpuri language localization support.", TEAL),
        ("Phase 4 (Q2 2027)", "AR Heritage Tours", "Integrate AR camera view for historical monument guides & virtual tours.", EMERALD)
    ]
    for idx, (pname, title, desc, col) in enumerate(phases):
        x = Inches(0.8 + idx * 2.95); y = Inches(1.8)
        card = s20.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
        hbar = s20.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(0.85))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = col; hbar.line.fill.background()
        tf_h = hbar.text_frame; tf_h.word_wrap = True; p_h = tf_h.paragraphs[0]; p_h.text = pname; p_h.font.size = Pt(12); p_h.font.bold = True; p_h.font.color.rgb = CARD_WHITE; p_h.alignment = PP_ALIGN.CENTER
        tb = s20.shapes.add_textbox(x + Inches(0.15), y + Inches(1.0), Inches(2.45), Inches(3.6))
        tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAVY
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(12); p2.font.color.rgb = TEXT_DARK; p2.space_before = Pt(8)
    add_notes(s20, "Looking ahead, our roadmap includes offline PWA tile caching, multi-lingual support, and AR-guided monument exploration.")

    # SLIDE 21: Conclusion
    s21 = create_base("Conclusion & Project Achievements", "CONCLUSION")
    c = s21.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.6))
    c.fill.solid(); c.fill.fore_color.rgb = CARD_WHITE; c.line.color.rgb = BORDER_LIGHT
    tb = s21.shapes.add_textbox(Inches(1.1), Inches(1.65), Inches(11.133), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = "🏆 Project Summary & Achievements"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = NAVY
    concs = [
        "Mission Accomplished: Successfully engineered a secure, GIS-enabled tourism exploration platform for Bihar.",
        "GIS Intelligence: Leaflet.js map engine with dynamic clustering, category pins, and straight-line distance calculations.",
        "AI Trip Planner: Haversine-based spatial clustering engine for automated 1-5 day itineraries with PDF export.",
        "Enterprise Governance: 11 Admin modules, soft-delete recycle bin, and immutable audit action logs.",
        "Production Quality: 87 unit tests passed with 100% pass rate, security hardened against OWASP vulnerabilities."
    ]
    for m in concs:
        p = tf.add_paragraph(); p.text = "• " + m; p.font.size = Pt(13); p.font.color.rgb = TEXT_DARK; p.space_before = Pt(12)
    add_notes(s21, "In conclusion, HiddenYatra delivers an impactful, production-ready solution combining modern web design with strong software engineering principles.")

    # SLIDE 22: Thank You
    s22 = prs.slides.add_slide(blank_layout); set_bg(s22, NAVY)
    l = s22.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    l.fill.solid(); l.fill.fore_color.rgb = SAFFRON; l.line.fill.background()
    card = s22.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card.fill.solid(); card.fill.fore_color.rgb = CARD_WHITE; card.line.color.rgb = BORDER_LIGHT
    tb = s22.shapes.add_textbox(Inches(1.8), Inches(1.5), Inches(9.733), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.text = "T H A N K   Y O U !"; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = SAFFRON; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Questions & Answers Session"; p2.font.size = Pt(24); p2.font.bold = True; p2.font.color.rgb = NAVY; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(8)
    p3 = tf.add_paragraph(); p3.text = "We welcome your questions, feedback, and technical evaluation."; p3.font.size = Pt(15); p3.font.color.rgb = TEXT_MUTED; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(20)
    p4 = tf.add_paragraph(); p4.text = "Project Codebase: d:\\HiddenYatra\nLive App Endpoint: http://127.0.0.1:5000/"; p4.font.size = Pt(14); p4.font.color.rgb = INDIGO; p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(20)
    add_notes(s22, "Thank you respected examiners for your time and guidance. We are ready for the Viva Voce examination.")

    # Save to workspace root AND artifacts
    prs.save(out_filename)
    print(f"Successfully generated PowerPoint deck ({len(prs.slides)} slides): {out_filename}")

if __name__ == "__main__":
    build_deck(r"d:\HiddenYatra\HiddenYatra_College_Major_Project_Presentation.pptx")
    build_deck(r"d:\HiddenYatra\HiddenYatra_College_Major_Project_Presentation_Final.pptx")
    build_deck(r"C:\Users\AKARSH RAJ\.gemini\antigravity-ide\brain\7922bebd-209c-4931-8ed1-a1e42ee92256\HiddenYatra_College_Major_Project_Presentation.pptx")
    build_deck(r"C:\Users\AKARSH RAJ\.gemini\antigravity-ide\brain\7922bebd-209c-4931-8ed1-a1e42ee92256\HiddenYatra_College_Major_Project_Presentation_Final.pptx")
